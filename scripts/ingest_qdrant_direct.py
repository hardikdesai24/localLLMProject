# ingest_qdrant_direct.py

# IMPORTS
import os
import time
import uuid
import io
from pathlib import Path

import pdfplumber
import pandas as pd
import pymupdf
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Users\Ananya.Mehta\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"
from PIL import Image
from pptx import Presentation
from docx import Document as DocxDocument

from openai import OpenAI
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

from llama_index.core import Document
from llama_index.core.node_parser import SentenceSplitter

# CONFIG
from dotenv import load_dotenv
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
EMBED_MODEL = os.getenv("EMBED_MODEL", "text-embedding-3-small")
QDRANT_HOST = os.getenv("QDRANT_HOST", "localhost")
QDRANT_PORT = int(os.getenv("QDRANT_PORT", "6333"))
DOCS_PATH = os.getenv("DOCS_PATH")

DIMENSIONS=int(os.getenv("DIMENSIONS", "512"))      # reduce from default 1536
BATCH_SIZE = int(os.getenv("BATCH_SIZE", "100"))      # points per Qdrant upsert
EMBED_BATCH = int(os.getenv("EMBED_BATCH", "100"))       # chunks per OpenAI API call

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100

SKIP_FILES = ["errors.log"]


# CLIENTS
openai_client = OpenAI(api_key=OPENAI_API_KEY)
qdrant_client = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT, timeout=300)

# HELPERS
def to_collection_name(filename: str) -> str:
    base = os.path.splitext(os.path.basename(filename))[0]
    base = base.lower().replace(" ", "_").replace("-", "_")
    return f"{base}_{DIMENSIONS}"

def embed_chunks(chunks):
    all_embeddings = []
    total = len(chunks)
    for i in range(0, total, EMBED_BATCH):
        batch = chunks[i : i + EMBED_BATCH]
        resp = openai_client.embeddings.create(
            model=EMBED_MODEL,
            input=batch,
            dimensions=DIMENSIONS,
        )
        all_embeddings.extend([r.embedding for r in resp.data])
        print(f"  -> Embedded {min(i + EMBED_BATCH, total)}/{total} chunks")
    return all_embeddings

def load_image_as_doc(path, file_name):
    img = Image.open(path).convert("RGB")
    text = pytesseract.image_to_string(img, config="--psm 6")
    return [Document(text=text, metadata={"file_name": file_name, "source_type": "image"})]
def load_pdf_as_single_doc(path, file_name):
    pages = []
    pdf_img_doc = pymupdf.open(path)

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            txt = page.extract_text() or ""
            if len(txt.strip()) < 20:
                try:
                    pymupdf_page = pdf_img_doc.load_page(i)
                    pix = pymupdf_page.get_pixmap(dpi=200)
                    img_bytes = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_bytes))
                    txt = pytesseract.image_to_string(img, config)
                except Exception as e:
                    print(f"  -> OCR failed on page {i+1}: {e}")
            pages.append(txt)

    pdf_img_doc.close()
    full_text = "\n\n".join(pages)
    return [Document(text=full_text, metadata={"file_name": file_name, "source_type": "pdf"})]

def load_csv_as_doc(path, file_name):
    df = pd.read_csv(path, encoding="utf-8", on_bad_lines="skip")
    text = df.fillna("").astype(str).to_csv(index=False)
    return [Document(text=text, metadata={"file_name": file_name, "source_type": "csv"})]


def load_excel_as_doc(path, file_name):
    xls = pd.ExcelFile(path, engine="openpyxl")
    parts = []
    for sheet in xls.sheet_names:
        df = xls.parse(sheet)
        if df.empty:
            continue
        csv_text = df.fillna("").astype(str).to_csv(index=False)
        parts.append(f"# Sheet: {sheet}\n\n{csv_text}")
    full_text = "\n\n".join(parts)
    return [Document(text=full_text, metadata={"file_name": file_name, "source_type": "excel"})]

def load_ppt_as_doc(path, file_name):

    prs = Presentation(path)
    docs = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        
        # TEXT EXTRACTION
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_content.append(text)
 
        # OCR FOR IMAGES
        for shape in slide.shapes:
            try:
                if shape.shape_type == 13:  
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                    ocr_text = pytesseract.image_to_string(img, config="--psm 6")
                    if ocr_text.strip():
                        slide_content.append(ocr_text.strip())
 
            except Exception as e:
                print(f"  -> OCR failed on slide {i+1}: {e}")
 
        # CREATE DOCUMENT PER SLIDE
        if slide_content:
            combined = "\n".join(slide_content)
            docs.append(
                Document(
                    text=combined,
                    metadata={"file_name": file_name,"source_type": "ppt","slide_number": i + 1}
                )
            )
    return docs

def load_docx_as_doc(path, file_name):
    doc = DocxDocument(path)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    full_text = "\n\n".join(paragraphs)
    return [Document(text=full_text, metadata={"file_name": file_name, "source_type": "docx"})]
 
def load_txt_as_doc(path, file_name):
    try:
        with open(path, "r", encoding="utf-8") as f:
            text = f.read()
    except UnicodeDecodeError:
        # fallback for weird encodings
        with open(path, "r", encoding="latin-1") as f:
            text = f.read()
 
    return [Document(text=text, metadata={"file_name": file_name, "source_type": "txt"})]

# QDRANT INGEST
def ingest_to_qdrant(collection_name, chunks, metadata_list):
    # Ghost collection cleanup
    if qdrant_client.collection_exists(collection_name):
        info = qdrant_client.get_collection(collection_name)
        if info.points_count == 0:
            qdrant_client.delete_collection(collection_name)
            print(f"  -> Cleared empty ghost: {collection_name}")
        else:
            print(f"  -> Collection already populated — skipping ingest")
            return

    # Create collection with 512-dim vectors
    qdrant_client.create_collection(
        collection_name=collection_name,
        vectors_config=VectorParams(size=DIMENSIONS, distance=Distance.COSINE),
    )

    print(f"  -> Generating embeddings ({len(chunks)} chunks)...")
    embeddings = embed_chunks(chunks)

    points = [
        PointStruct(
            id=str(uuid.uuid4()),
            vector=embeddings[i],
            payload={
                **metadata_list[i],
                "text": chunks[i],         
            },
        )
        for i in range(len(chunks))
    ]

    total = len(points)
    for i in range(0, total, BATCH_SIZE):
        batch = points[i : i + BATCH_SIZE]
        qdrant_client.upsert(
            collection_name=collection_name,
            points=batch,
            wait=True,
        )
        written = min(i + BATCH_SIZE, total)
        print(f"  -> Written {written}/{total} points to Qdrant")
        time.sleep(0.05)

    print(f" ✅ Done — {total} points in [{collection_name}]\n")

# MAIN
if __name__ == "__main__":
    if not os.path.isdir(DOCS_PATH):
        print(f"\n[ERROR] DOCS_PATH folder not found: {DOCS_PATH}")
        raise SystemExit(1)
    
    files = [f for f in os.listdir(DOCS_PATH) if f not in SKIP_FILES and os.path.isfile(os.path.join(DOCS_PATH, f))]

    if not files:
        print(f"\n[WARN] No files found in {DOCS_PATH}")
        raise SystemExit(0)

    print("\n[INIT] Loading documents and configuring ingestion...")
    print(f"       Folder    : {DOCS_PATH}")
    print(f"       Files     : {len(files)}")
    print(f"       Embedding : {EMBED_MODEL} (dims={DIMENSIONS})")

    total_files = len(files)
    completed = 0
    failed = 0

    print("\n[START] Ingesting files into Qdrant...\n")

    for idx, fname in enumerate(files, start=1):
        filepath   = os.path.join(DOCS_PATH, fname)
        collection = to_collection_name(fname)

        print("============================================================")
        print(f"[{idx}/{total_files}] {fname}")
        try:
            size_mb = os.path.getsize(filepath) / (1024 * 1024)
            print(f"        → Size       : {size_mb:.1f} MB")
            print(f"        → Collection : {collection}")

            ext = Path(fname).suffix.lower()

            if ext == ".pdf":
                docs = load_pdf_as_single_doc(filepath, fname)
            elif ext == ".csv":
                docs = load_csv_as_doc(filepath, fname)
            elif ext in [".xlsx", ".xls"]:
                docs = load_excel_as_doc(filepath, fname)
            elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif"]:
                docs = load_image_as_doc(filepath, fname)
            elif ext == ".pptx":
                docs = load_ppt_as_doc(filepath, fname)
            elif ext == ".docx":
                docs = load_docx_as_doc(filepath, fname)
            elif ext == ".txt":
                docs = load_txt_as_doc(filepath, fname)
            else:
                print(f"        → Skipping unsupported file type: {ext}")
                continue

            parser = SentenceSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
            nodes = parser.get_nodes_from_documents(docs)
            chunks = [n.text for n in nodes]
            metadata_list = [n.metadata for n in nodes]



            print(f"        → Chunks     : {len(chunks)}")
            print(f"        → Embedding and indexing...")

            ingest_to_qdrant(collection, chunks, metadata_list)
            completed += 1

        except Exception as e:
            print(f"\n        ❌ Failed: {e}\n")
            failed += 1

    print("============================================================")
    print("\n✅ Ingestion Summary:")
    print(f"   Completed : {completed} collections")
    print(f"   Skipped   : 0 files")
    print(f"   Failed    : {failed} files")
    print("   Folder    :", DOCS_PATH)
    print("\nVerify at: http://localhost:6333/dashboard")